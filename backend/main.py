from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pptx import Presentation
from pydantic import BaseModel
from openai import OpenAI
from dotenv import load_dotenv
import io
import os
import logging

# Load environment variables
load_dotenv()

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:5174", "http://localhost:5175"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Only mount static files if the directory exists (production mode)
frontend_dist = "../frontend/dist"
if os.path.exists(frontend_dist):
    app.mount("/", StaticFiles(directory=frontend_dist, html=True), name="static")
    if os.path.exists(os.path.join(frontend_dist, "assets")):
        app.mount("/assets", StaticFiles(directory=os.path.join(frontend_dist, "assets")), name="assets")

# Serve the React app at the root
@app.get("/{full_path:path}")
async def serve_spa(full_path: str):
    if full_path.startswith("api/"):
        raise HTTPException(status_code=404, detail="Not found")
    return FileResponse("../frontend/dist/index.html")

class FeedbackRequest(BaseModel):
    titles: str
    targetAudience: str

def get_openai_client():
    """Initialize OpenAI client only when needed"""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OpenAI API key not configured")
    return OpenAI()  # It will automatically use OPENAI_API_KEY from environment

@app.post("/api/get-feedback")
async def get_feedback(request: FeedbackRequest):
    try:
        try:
            client = get_openai_client()
        except ValueError as e:
            logger.error("OpenAI API key is not configured")
            raise HTTPException(status_code=500, detail="OpenAI API key not configured. Please check server configuration.")

        logger.info(f"Generating feedback for audience: {request.targetAudience}")
        logger.info(f"Content length: {len(request.titles)} characters")

        prompt = f"""This narrative is intended for {request.targetAudience}.

Content to analyze:
{request.titles}

Please provide feedback on the following:
1. Check if the introduction effectively frames the problem and if the conclusion clearly delivers a strong call to action or takeaway.
2. Identify any points where the transition between slides feels abrupt, confusing, or could be improved.
3. Suggest how to reorganize the slides to make the argument more persuasive."""

        try:
            logger.info("Sending request to OpenAI API...")
            completion = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a presentation structure expert. Analyze the slide titles and provide constructive feedback on the narrative flow."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=1000
            )
            logger.info("Successfully received response from OpenAI")
            feedback = completion.choices[0].message.content
            return {"feedback": feedback}
        except Exception as e:
            logger.error(f"OpenAI API error details: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error getting feedback from AI service: {str(e)}"
            )

    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Unexpected error in get_feedback: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Server error while processing feedback request: {str(e)}"
        )

# API endpoints under /api prefix
@app.post("/api/extract-titles")
async def extract_titles(file: UploadFile = File(...)):
    try:
        # Log request received
        logger.info(f"Received file: {file.filename}")
        
        # Verify file type
        if not file.filename.endswith('.pptx'):
            logger.error(f"Invalid file type: {file.filename}")
            raise HTTPException(status_code=400, detail="File must be a .pptx file")
        
        # Read the uploaded file
        try:
            content = await file.read()
            if not content:
                logger.error("Empty file received")
                raise HTTPException(status_code=400, detail="Empty file uploaded")
        except Exception as e:
            logger.error(f"Error reading file: {str(e)}")
            raise HTTPException(status_code=400, detail="Could not read the uploaded file")
            
        # Create presentation object
        try:
            pptx = Presentation(io.BytesIO(content))
        except Exception as e:
            logger.error(f"Error creating Presentation object: {str(e)}")
            raise HTTPException(status_code=400, detail="Invalid PowerPoint file format")
        
        # Extract titles and identify sections
        titles = []
        current_section = None
        
        for i, slide in enumerate(pptx.slides):
            try:
                # Get all text from the slide
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text.strip() + "\n"
                
                title_text = slide.shapes.title.text.strip() if slide.shapes.title else ""
                
                # Check if this might be a section slide
                is_section_slide = False
                if not title_text and slide_text.strip():
                    # If there's no title but there is other text, this might be a section slide
                    lines = [line.strip() for line in slide_text.split('\n') if line.strip()]
                    if lines:
                        current_section = lines[0]  # Use the first non-empty line as section title
                        is_section_slide = True
                        logger.info(f"Found section slide {i+1}: {current_section}")
                
                if is_section_slide:
                    # For section slides, use a special format
                    titles.append(f"[SECTION] {current_section}")
                elif title_text:
                    # For regular slides with titles
                    titles.append(title_text)
                    logger.info(f"Extracted title from slide {i+1}: {title_text}")
                else:
                    # For slides without titles
                    titles.append("[No Title]")
                    logger.info(f"No title found in slide {i+1}")
                
            except Exception as e:
                logger.error(f"Error processing slide {i+1}: {str(e)}")
                titles.append(f"[Error processing slide {i+1}]")
        
        if not titles:
            logger.warning("No slides found in presentation")
            return {"titles": "No slides found in the presentation"}
        
        # Create formatted text with titles
        formatted_titles = []
        current_section = None
        
        for i, title in enumerate(titles):
            if title.startswith("[SECTION]"):
                # Handle section markers
                current_section = title.replace("[SECTION]", "").strip()
                formatted_titles.append(f"\n{current_section}\n{'='*len(current_section)}")
            else:
                # Format regular slides
                prefix = f"p.{i+1}"
                if current_section:
                    formatted_titles.append(f"{prefix}: {title}")
                else:
                    formatted_titles.append(f"{prefix}: {title}")
        
        titles_text = "\n".join(formatted_titles)
        
        logger.info(f"Successfully processed {len(titles)} slides")
        return {"titles": titles_text}
        
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e)) 